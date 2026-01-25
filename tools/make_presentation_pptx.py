"""tools/make_presentation_pptx.py

Run:
    python tools/make_presentation_pptx.py
    python tools/make_presentation_pptx.py --out presentation/ZegarBiologiczny.pptx

Opis:
    Generuje profesjonalną prezentację .pptx dla projektu ZegarBiologiczny:
    - 16:9, spójny styl (kolor akcentu, stopka, numeracja),
    - slajdy sekcyjne (Agenda / Dane / Modele / Deploy),
    - prosty diagram pipeline,
    - metryki wciągane automatycznie z logów w Logs/ (jeśli dostępne).

Wymagania:
    pip install python-pptx
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from pathlib import Path
import argparse
import re
from typing import Optional


@dataclass(frozen=True)
class Slide:
    kind: str  # title | section | bullets | diagram
    title: str
    bullets: list[str] | None = None
    notes: str = ""


@dataclass(frozen=True)
class Metrics:
    baseline_rgb_acc: Optional[float] = None
    advanced_rf_acc: Optional[float] = None
    mlp_test_cmae_h: Optional[float] = None
    mlp_test_p90_h: Optional[float] = None
    mlp_test_p95_h: Optional[float] = None


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--out",
        type=Path,
        default=Path("presentation/ZegarBiologiczny.pptx"),
        help="Ścieżka wyjściowa .pptx",
    )
    parser.add_argument(
        "--no-log-metrics",
        action="store_true",
        help="Nie parsuj metryk z Logs/ (użyj wartości domyślnych)",
    )
    return parser.parse_args()


def _read_text_if_exists(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def _parse_first_float(pattern: str, text: str) -> Optional[float]:
    m = re.search(pattern, text, flags=re.IGNORECASE | re.MULTILINE)
    if not m:
        return None
    try:
        return float(m.group(1))
    except Exception:
        return None


def _find_latest_log(glob_pattern: str) -> Optional[Path]:
    logs_dir = Path("Logs")
    if not logs_dir.exists():
        return None
    candidates = list(logs_dir.glob(glob_pattern))
    if not candidates:
        return None
    candidates.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    return candidates[0]


def load_metrics_from_logs() -> Metrics:
    # baseline RGB
    rgb_log = Path("Logs/baseline_rgb_log.txt")
    rgb_text = _read_text_if_exists(rgb_log)
    baseline_rgb_acc = _parse_first_float(r"^Accuracy:\s*([0-9.]+)", rgb_text)

    # baseline advanced (RF)
    adv_log = Path("Logs/baseline_advanced_log.txt")
    adv_text = _read_text_if_exists(adv_log)
    # prefer explicit RF line
    advanced_rf_acc = _parse_first_float(r"^Accuracy \(rf\):\s*([0-9.]+)", adv_text)
    if advanced_rf_acc is None:
        advanced_rf_acc = _parse_first_float(
            r"Najlepszy model:\s*rf\s*\(accuracy\s*=\s*([0-9.]+)\)", adv_text
        )

    # MLP cyclic results (use latest by mtime)
    mlp_log = _find_latest_log("**/train_hour_nn_cyclic_*_log.txt")
    mlp_text = _read_text_if_exists(mlp_log) if mlp_log else ""
    mlp_test_cmae_h = _parse_first_float(r"^\[INFO\]\s*TEST\s*Cyclic MAE:\s*([0-9.]+)h", mlp_text)
    mlp_test_p90_h = _parse_first_float(r"^\[INFO\]\s*TEST\s*P90:\s*([0-9.]+)h", mlp_text)
    mlp_test_p95_h = _parse_first_float(r"^\[INFO\]\s*TEST\s*P95:\s*([0-9.]+)h", mlp_text)

    return Metrics(
        baseline_rgb_acc=baseline_rgb_acc,
        advanced_rf_acc=advanced_rf_acc,
        mlp_test_cmae_h=mlp_test_cmae_h,
        mlp_test_p90_h=mlp_test_p90_h,
        mlp_test_p95_h=mlp_test_p95_h,
    )


def build_slides(metrics: Metrics) -> list[Slide]:
    today = date.today().isoformat()

    rgb_acc = (
        f"Accuracy: {metrics.baseline_rgb_acc:.3f}" if metrics.baseline_rgb_acc is not None else "Accuracy: ≈ 0.24"
    )
    rf_acc = (
        f"Accuracy (RF): {metrics.advanced_rf_acc:.3f}" if metrics.advanced_rf_acc is not None else "Accuracy (RF): ≈ 0.79"
    )
    if metrics.mlp_test_cmae_h is not None:
        cmae = f"TEST Cyclic MAE: {metrics.mlp_test_cmae_h:.3f}h (~{metrics.mlp_test_cmae_h*60:.0f} min)"
    else:
        cmae = "TEST Cyclic MAE: ≈ 0.54h (~33 min)"
    p90 = f"P90: {metrics.mlp_test_p90_h:.3f}h" if metrics.mlp_test_p90_h is not None else "P90: ≈ 1.46h"
    p95 = f"P95: {metrics.mlp_test_p95_h:.3f}h" if metrics.mlp_test_p95_h is not None else "P95: ≈ 2.15h"

    return [
        Slide(
            kind="title",
            title="ZegarBiologiczny",
            bullets=[
                "Predykcja godziny (0–23) na podstawie zdjęć nieba",
                "ML pipeline: data → cechy → modele → Raspberry Pi overlay",
                f"Data: {today}",
            ],
            notes="1–2 zdania: problem → podejście → co pokażesz (pipeline + wyniki + deploy).",
        ),
        Slide(
            kind="bullets",
            title="Agenda",
            bullets=[
                "Motywacja i definicja problemu",
                "Dane + pipeline (walidacja, cechy, normalizacja)",
                "Modele i wyniki (baseline → robust/MLP)",
                "Deploy na Raspberry Pi + synchronizacja danych",
                "Wnioski i dalsze kroki",
            ],
        ),
        Slide(kind="section", title="Dane i pipeline"),
        Slide(
            kind="bullets",
            title="Problem i cel",
            bullets=[
                "Wejście: obraz z kamery (scena zewnętrzna / niebo)",
                "Wyjście: godzina dnia (0–23) lub czas cykliczny (sin/cos)",
                "Wyzwania: chmury, ekspozycja, sezonowość, różne kamery",
            ],
        ),
        Slide(
            kind="bullets",
            title="Dane i etykiety",
            bullets=[
                "Zbieranie: MLDailyHourClock.py",
                "Struktura: dataset/YYYY/MM/DD/HH/*.jpg",
                "labels.csv: filepath, hour, datetime",
                "Walidacja/ładowanie: src/load_data.py (single source of truth)",
            ],
        ),
        Slide(
            kind="diagram",
            title="Pipeline end-to-end (skrót)",
            bullets=[
                "Uruchomienie całości: ./run_full_pipeline.sh",
                "Bez data leakage: normalizacja liczona tylko na train",
            ],
        ),
        Slide(
            kind="bullets",
            title="Cechy: od prostych do robust",
            bullets=[
                "Mean RGB: szybki baseline (3 liczby)",
                "Advanced: RGB+HSV statystyki",
                "Robust: histogramy/entropia/gradient/edges → odporność na chmury",
            ],
        ),
        Slide(kind="section", title="Modele i wyniki"),
        Slide(
            kind="bullets",
            title="Modele (przegląd)",
            bullets=[
                "Baseline RGB: LogisticRegression",
                "Advanced: porównanie kilku klasyków (RF/GB/KNN/LogReg)",
                "Robust: regresja cykliczna (sin/cos) + MLP",
                "CNN: eksperymentalnie end-to-end",
            ],
        ),
        Slide(
            kind="bullets",
            title="Wyniki: baseline RGB (punkt odniesienia)",
            bullets=[
                rgb_acc,
                "Wniosek: kolor sceny daje ograniczony sygnał (benchmark)",
                "Źródło: Logs/baseline_rgb_log.txt",
            ],
        ),
        Slide(
            kind="bullets",
            title="Wyniki: advanced (RF)",
            bullets=[
                rf_acc,
                "Wniosek: ręcznie zaprojektowane cechy + klasyk ML działają dobrze",
                "Źródło: Logs/baseline_advanced_log.txt",
            ],
        ),
        Slide(
            kind="bullets",
            title="Wyniki: MLP cykliczny (robust)",
            bullets=[
                "Kodowanie czasu: (sin θ, cos θ), θ = 2πh/24",
                cmae,
                f"{p90}, {p95}",
                "Źródło: Logs/**/train_hour_nn_cyclic_*_log.txt",
            ],
        ),
        Slide(kind="section", title="Deploy i operacje"),
        Slide(
            kind="bullets",
            title="Deploy: Raspberry Pi overlay",
            bullets=[
                "Desktop: src/camera_hour_overlay*.py",
                "RPi: src/camera_hour_overlay_mlp_rpi.py (MLP lub fallback RF)",
                "Artefakty: models/pc/ (trening) → models/rpi/ (deploy)",
            ],
        ),
        Slide(
            kind="bullets",
            title="Synchronizacja danych RPi → PC/Windows",
            bullets=[
                "Worker: synchro_dataset.sh (SRC → STAGING → rsync)",
                "Retry + logowanie + opcjonalny WoL",
                "Optymalizacja skanowania przy dużym dataset: SCAN_MODE=recent-hours",
            ],
        ),
        Slide(
            kind="bullets",
            title="Wnioski i kolejne kroki",
            bullets=[
                "Ustabilizować ewaluację (time-based split jako standard)",
                "Walidacja na innych kamerach/lokalizacjach (generalizacja)",
                "Monitoring driftu (pogoda/sezon) i strategia retrain",
                "Ewentualnie: transfer learning + augmentacje",
            ],
        ),
    ]


def _ensure_dep() -> None:
    try:
        import pptx  # noqa: F401
    except Exception as exc:  # pragma: no cover
        raise SystemExit(
            "Brak zależności 'python-pptx'. Zainstaluj: pip install python-pptx"
        ) from exc


def _add_bullets(text_frame, bullets: list[str]) -> None:
    # pierwsza linia jest w domyślnym paragrafie
    if not bullets:
        return

    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = bullets[0]
    p.level = 0

    for b in bullets[1:]:
        p = text_frame.add_paragraph()
        p.text = b
        p.level = 0

def main() -> None:
    args = parse_args()
    _ensure_dep()

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    out_path: Path = args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    metrics = Metrics()
    if not args.no_log_metrics:
        metrics = load_metrics_from_logs()

    prs = Presentation()
    # 16:9 wide
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    accent = RGBColor(0x1F, 0x4E, 0x79)  # stonowany niebieski
    dark = RGBColor(0x1A, 0x1A, 0x1A)
    light = RGBColor(0xFF, 0xFF, 0xFF)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]

    slides = build_slides(metrics)
    total = len(slides)

    def style_title(shape) -> None:
        tf = shape.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(36)
                r.font.color.rgb = dark

    def style_body(shape) -> None:
        tf = shape.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Calibri"
                r.font.size = Pt(22)
                r.font.color.rgb = dark

    def add_header_bar(slide) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

    def add_footer(slide, idx: int) -> None:
        left = Inches(0.5)
        top = prs.slide_height - Inches(0.45)
        width = prs.slide_width - Inches(1.0)
        height = Inches(0.3)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"ZegarBiologiczny  |  {idx}/{total}"
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    def set_white_bg(slide) -> None:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = light

    def add_pipeline_diagram(slide) -> None:
        # Prosty diagram: 5 boksów + strzałki
        y = Inches(3.1)
        x0 = Inches(0.7)
        w = Inches(2.2)
        h = Inches(0.7)
        gap = Inches(0.35)
        labels = ["Zbieranie danych", "Etykiety", "Cechy", "Normalizacja", "Trening/Deploy"]

        shapes = []
        for i, lab in enumerate(labels):
            x = x0 + i * (w + gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFB)
            box.line.color.rgb = accent
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = lab
            p.font.name = "Calibri"
            p.font.size = Pt(16)
            p.font.color.rgb = dark
            shapes.append(box)

        # connectors as thin rectangles (prostsze niż konektory w python-pptx)
        for i in range(len(shapes) - 1):
            x = shapes[i].left + shapes[i].width
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x, y + Inches(0.18), gap, Inches(0.34)
            )
            conn.fill.solid()
            conn.fill.fore_color.rgb = accent
            conn.line.fill.background()

    for idx, s in enumerate(slides, start=1):
        if s.kind == "title":
            slide = prs.slides.add_slide(title_layout)
            set_white_bg(slide)
            add_header_bar(slide)
            slide.shapes.title.text = s.title
            style_title(slide.shapes.title)

            # subtitle placeholder
            subtitle = slide.placeholders[1]
            subtitle.text = "\n".join(s.bullets or [])
            for p in subtitle.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Calibri"
                    r.font.size = Pt(20)
                    r.font.color.rgb = dark
            add_footer(slide, idx)
            if s.notes:
                slide.notes_slide.notes_text_frame.text = s.notes
            continue

        if s.kind == "section":
            slide = prs.slides.add_slide(blank_layout)
            set_white_bg(slide)
            add_header_bar(slide)

            # duży tytuł sekcji
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), prs.slide_width - Inches(1.6), Inches(1.2))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = s.title
            p.font.name = "Calibri"
            p.font.size = Pt(44)
            p.font.color.rgb = accent
            add_footer(slide, idx)
            continue

        if s.kind == "diagram":
            slide = prs.slides.add_slide(blank_layout)
            set_white_bg(slide)
            add_header_bar(slide)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), prs.slide_width - Inches(1.6), Inches(0.6))
            title_tf = title_box.text_frame
            title_tf.clear()
            title_tf.paragraphs[0].text = s.title
            title_tf.paragraphs[0].font.name = "Calibri"
            title_tf.paragraphs[0].font.size = Pt(32)
            title_tf.paragraphs[0].font.color.rgb = dark
            add_pipeline_diagram(slide)
            # krótkie bullet points pod diagramem
            if s.bullets:
                tb = slide.shapes.add_textbox(Inches(0.9), Inches(4.2), prs.slide_width - Inches(1.8), Inches(2.5))
                _add_bullets(tb.text_frame, s.bullets)
                style_body(tb)
            add_footer(slide, idx)
            if s.notes:
                slide.notes_slide.notes_text_frame.text = s.notes
            continue

        # default: bullets
        slide = prs.slides.add_slide(content_layout)
        set_white_bg(slide)
        add_header_bar(slide)
        slide.shapes.title.text = s.title
        style_title(slide.shapes.title)

        # standard content placeholder
        body = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                body = shape
                break

        if body is not None and s.bullets:
            _add_bullets(body.text_frame, s.bullets)
            style_body(body)
        add_footer(slide, idx)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes

    prs.save(out_path)
    print(f"[OK] Zapisano: {out_path}")


if __name__ == "__main__":
    main()
